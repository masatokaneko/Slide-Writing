from pathlib import Path
import os
from datetime import datetime
from pptx import Presentation
from pptx.util import Cm, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import tempfile

class PPTXCreator:
    # BCGカラーパレット
    BCG_BLUE = RGBColor(0, 112, 192)
    SECONDARY_BLUE = RGBColor(68, 114, 196)
    ACCENT_ORANGE = RGBColor(255, 192, 0)
    DARK_GRAY = RGBColor(64, 64, 64)
    LIGHT_GRAY = RGBColor(217, 217, 217)
    WHITE = RGBColor(255, 255, 255)

    # フォント階層
    FONT_TITLE = ("游ゴシック", 36, True)
    FONT_SUBTITLE = ("游ゴシック", 24, False)
    FONT_HEADER = ("游ゴシック", 28, True)
    FONT_BODY = ("游ゴシック", 18, False)
    FONT_CAPTION = ("游ゴシック", 14, False)

    def __init__(self, output_dir=None):
        if output_dir is None:
            self.output_dir = Path("data") / "generated"
        else:
            self.output_dir = Path(output_dir)
        self.output_dir.mkdir(parents=True, exist_ok=True)
        self.fonts = ["游ゴシック", "Yu Gothic", "メイリオ", "Meiryo", "sans-serif"]
        self.slide_width = Cm(33.867)
        self.slide_height = Cm(19.05)

    def create_presentation(self, structure: dict) -> str:
        """
        スライド構造からPowerPointファイルを生成し、パスを返す
        """
        prs = Presentation()
        prs.slide_width = self.slide_width
        prs.slide_height = self.slide_height
        try:
            slides = structure.get("slides", [])
            for slide in slides:
                stype = slide.get("type", "content_slide")
                if stype == "title_slide":
                    self._add_title_slide(prs, slide, structure.get("title", ""))
                elif stype == "content_slide":
                    self._add_content_slide(prs, slide)
                elif stype == "financial_slide":
                    self._add_financial_slide(prs, slide)
                elif stype == "implementation_slide":
                    self._add_implementation_slide(prs, slide)
                else:
                    self._add_content_slide(prs, slide)
            now = datetime.now().strftime("%Y%m%d_%H%M%S")
            filename = f"presentation_{now}.pptx"
            output_path = self.output_dir / filename
            with tempfile.NamedTemporaryFile(delete=False, suffix=".pptx") as tmp:
                tmp_path = Path(tmp.name)
                try:
                    prs.save(str(tmp_path))
                    tmp.close()
                    tmp_path.replace(output_path)
                except Exception as e:
                    return f"ファイル保存エラー: {str(e)}"
                finally:
                    if tmp_path.exists():
                        try:
                            tmp_path.unlink()
                        except Exception:
                            pass
            return str(output_path)
        except Exception as e:
            return f"PPTX生成エラー: {str(e)}"

    def _set_font(self, shape, size=18, bold=False, color=None):
        for p in shape.text_frame.paragraphs:
            for run in p.runs:
                run.font.size = Pt(size)
                run.font.bold = bold
                run.font.name = self.fonts[0]
                if color:
                    run.font.color.rgb = color

    def _add_title_slide(self, prs, slide, pres_title):
        layout = prs.slide_layouts[6]  # 空白
        s = prs.slides.add_slide(layout)
        # 背景グラデーション
        fill = s.background.fill
        fill.gradient()
        fill.gradient_stops[0].color.rgb = self.BCG_BLUE
        fill.gradient_stops[1].color.rgb = self.SECONDARY_BLUE
        # タイトル
        title = slide.get("title", pres_title)
        title_shape = s.shapes.add_textbox(Cm(2), Cm(6), Cm(30), Cm(4))
        title_shape.text = title
        self._set_font(title_shape, size=self.FONT_TITLE[1], bold=self.FONT_TITLE[2], color=self.WHITE)
        for p in title_shape.text_frame.paragraphs:
            p.alignment = PP_ALIGN.CENTER
        # サブタイトル
        subtitle = slide["content"].get("main_message", "")
        if subtitle:
            sub_shape = s.shapes.add_textbox(Cm(2), Cm(10), Cm(30), Cm(2.5))
            sub_shape.text = subtitle
            self._set_font(sub_shape, size=self.FONT_SUBTITLE[1], color=self.WHITE)
            for p in sub_shape.text_frame.paragraphs:
                p.alignment = PP_ALIGN.CENTER
        # 左下ライン装飾
        line = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(2), Cm(17), Cm(10), Cm(0.5))
        line.fill.solid()
        line.fill.fore_color.rgb = self.BCG_BLUE
        line.line.color.rgb = self.BCG_BLUE
        line.shadow.inherit = False

    def _add_content_slide(self, prs, slide):
        layout = prs.slide_layouts[6]
        s = prs.slides.add_slide(layout)
        # 背景色
        fill = s.background.fill
        fill.solid()
        fill.fore_color.rgb = self.WHITE
        # ヘッダーライン
        header = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(0), Cm(0), self.slide_width, Cm(0.5))
        header.fill.solid()
        header.fill.fore_color.rgb = self.BCG_BLUE
        header.line.color.rgb = self.BCG_BLUE
        # タイトル
        title = slide.get("title", "")
        title_shape = s.shapes.add_textbox(Cm(2), Cm(1), Cm(29), Cm(2.5))
        title_shape.text = title
        self._set_font(title_shape, size=self.FONT_HEADER[1], bold=self.FONT_HEADER[2], color=self.DARK_GRAY)
        # 本文（メインメッセージ）
        main_msg = slide["content"].get("main_message", "")
        msg_shape = s.shapes.add_textbox(Cm(2), Cm(4), Cm(20), Cm(3))
        msg_shape.text = main_msg
        self._set_font(msg_shape, size=self.FONT_BODY[1], color=self.DARK_GRAY)
        # 箇条書き
        points = slide["content"].get("supporting_points", [])[:5]
        y = 7
        for pt in points:
            box = s.shapes.add_textbox(Cm(3), Cm(y), Cm(25), Cm(1.5))
            box.text = pt
            self._set_font(box, size=self.FONT_BODY[1], color=self.DARK_GRAY)
            for p in box.text_frame.paragraphs:
                p.alignment = PP_ALIGN.LEFT
            y += 2
        # 右側メトリクス
        data = slide["content"].get("data", {})
        if data:
            metrics_box = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(24), Cm(4), Cm(7), Cm(6))
            metrics_box.fill.solid()
            metrics_box.fill.fore_color.rgb = self.LIGHT_GRAY
            metrics_box.line.color.rgb = self.SECONDARY_BLUE
            metrics_box.shadow.inherit = True
            tf = metrics_box.text_frame
            tf.text = "メトリクス"
            for k, v in data.items():
                p = tf.add_paragraph()
                p.text = f"{k}: {v}"
                p.font.size = Pt(self.FONT_CAPTION[1])
                p.font.name = self.FONT_CAPTION[0]

    def _add_financial_slide(self, prs, slide):
        layout = prs.slide_layouts[6]
        s = prs.slides.add_slide(layout)
        fill = s.background.fill
        fill.solid()
        fill.fore_color.rgb = self.WHITE
        # ヘッダーライン
        header = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(0), Cm(0), self.slide_width, Cm(0.5))
        header.fill.solid()
        header.fill.fore_color.rgb = self.ACCENT_ORANGE
        header.line.color.rgb = self.ACCENT_ORANGE
        # タイトル
        title = slide.get("title", "")
        title_shape = s.shapes.add_textbox(Cm(2), Cm(1), Cm(29), Cm(2.5))
        title_shape.text = title
        self._set_font(title_shape, size=self.FONT_HEADER[1], bold=self.FONT_HEADER[2], color=self.DARK_GRAY)
        # メイン数値ハイライト
        data = slide["content"].get("data", {})
        y = 4
        for k, v in data.items():
            num_box = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(3), Cm(y), Cm(10), Cm(2))
            num_box.fill.solid()
            num_box.fill.fore_color.rgb = self.LIGHT_GRAY
            num_box.line.color.rgb = self.ACCENT_ORANGE
            num_box.shadow.inherit = True
            tf = num_box.text_frame
            tf.text = f"{k}: {v}"
            tf.paragraphs[0].font.size = Pt(self.FONT_BODY[1])
            tf.paragraphs[0].font.bold = True
            tf.paragraphs[0].font.name = self.FONT_BODY[0]
            y += 2.5
        # 本文
        main_msg = slide["content"].get("main_message", "")
        msg_shape = s.shapes.add_textbox(Cm(15), Cm(4), Cm(15), Cm(6))
        msg_shape.text = main_msg
        self._set_font(msg_shape, size=self.FONT_BODY[1], color=self.DARK_GRAY)
        # データソース
        caption = s.shapes.add_textbox(Cm(2), Cm(16.5), Cm(25), Cm(1))
        caption.text = "出典: 社内データ・外部調査等"
        self._set_font(caption, size=self.FONT_CAPTION[1], color=self.DARK_GRAY)

    def _add_implementation_slide(self, prs, slide):
        layout = prs.slide_layouts[6]
        s = prs.slides.add_slide(layout)
        fill = s.background.fill
        fill.solid()
        fill.fore_color.rgb = self.WHITE
        # ヘッダーライン
        header = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(0), Cm(0), self.slide_width, Cm(0.5))
        header.fill.solid()
        header.fill.fore_color.rgb = self.SECONDARY_BLUE
        header.line.color.rgb = self.SECONDARY_BLUE
        # タイトル
        title = slide.get("title", "")
        title_shape = s.shapes.add_textbox(Cm(2), Cm(1), Cm(29), Cm(2.5))
        title_shape.text = title
        self._set_font(title_shape, size=self.FONT_HEADER[1], bold=self.FONT_HEADER[2], color=self.DARK_GRAY)
        # タイムライン（簡易）
        y = 4
        points = slide["content"].get("supporting_points", [])[:5]
        for i, pt in enumerate(points):
            phase_box = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(3), Cm(y), Cm(25), Cm(1.5))
            phase_box.fill.solid()
            phase_box.fill.fore_color.rgb = self.LIGHT_GRAY if i % 2 == 0 else self.WHITE
            phase_box.line.color.rgb = self.SECONDARY_BLUE
            phase_box.shadow.inherit = True
            tf = phase_box.text_frame
            tf.text = pt
            tf.paragraphs[0].font.size = Pt(self.FONT_BODY[1])
            tf.paragraphs[0].font.name = self.FONT_BODY[0]
            y += 2
        # アクション項目
        main_msg = slide["content"].get("main_message", "")
        action_shape = s.shapes.add_textbox(Cm(3), Cm(15), Cm(25), Cm(2))
        action_shape.text = main_msg
        self._set_font(action_shape, size=self.FONT_BODY[1], color=self.ACCENT_ORANGE)
        # 責任者・期限（仮）
        caption = s.shapes.add_textbox(Cm(2), Cm(17.5), Cm(25), Cm(1))
        caption.text = "責任者・期限: 詳細は別紙参照"
        self._set_font(caption, size=self.FONT_CAPTION[1], color=self.DARK_GRAY)

# 追加: main.pyから直接呼び出せる関数
def create_presentation(structure: dict, output_path: str = None) -> str:
    creator = PPTXCreator()
    result = creator.create_presentation(structure)
    if output_path and isinstance(result, str) and not result.startswith("PPTX生成エラー"):
        # ファイルを指定パスに移動
        try:
            Path(result).replace(output_path)
            return output_path
        except Exception as e:
            return f"ファイル移動エラー: {str(e)}"
    return result 