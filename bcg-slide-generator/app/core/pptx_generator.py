from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os
import json
from datetime import datetime

class PPTXGenerator:
    def __init__(self):
        self.slide_width = Inches(13.33)
        self.slide_height = Inches(7.5)
        self.margin = Inches(0.5)
        self.title_font_size = Pt(44)
        self.subtitle_font_size = Pt(32)
        self.body_font_size = Pt(24)
        self.bullet_font_size = Pt(20)
    
    def create_presentation(self, structure, design_pattern):
        """プレゼンテーションを生成"""
        prs = Presentation()
        
        # スライドサイズの設定
        prs.slide_width = self.slide_width
        prs.slide_height = self.slide_height
        
        # 各スライドの生成
        for slide_data in structure['slides']:
            slide_type = slide_data['type']
            content = slide_data['content']
            
            if slide_type == 'title_slide':
                self._create_title_slide(prs, content, design_pattern)
            elif slide_type == 'content_slide':
                self._create_content_slide(prs, content, design_pattern)
            elif slide_type == 'solution_slide':
                self._create_solution_slide(prs, content, design_pattern)
            elif slide_type == 'chart_slide':
                self._create_chart_slide(prs, content, design_pattern)
            elif slide_type == 'financial_slide':
                self._create_financial_slide(prs, content, design_pattern)
            elif slide_type == 'conclusion_slide':
                self._create_conclusion_slide(prs, content, design_pattern)
        
        return prs
    
    def _create_title_slide(self, prs, content, design_pattern):
        """タイトルスライドの生成"""
        slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)
        
        # タイトル
        title = slide.shapes.title
        title.text = content['main_message']
        title.text_frame.paragraphs[0].font.size = self.title_font_size
        title.text_frame.paragraphs[0].font.bold = True
        title.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # サブタイトル
        if 'subtitle' in content:
            subtitle = slide.placeholders[1]
            subtitle.text = content['subtitle']
            subtitle.text_frame.paragraphs[0].font.size = self.subtitle_font_size
            subtitle.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    def _create_content_slide(self, prs, content, design_pattern):
        """コンテンツスライドの生成"""
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        
        # タイトル
        title = slide.shapes.title
        title.text = content['main_message']
        title.text_frame.paragraphs[0].font.size = self.body_font_size
        title.text_frame.paragraphs[0].font.bold = True
        
        # 本文
        body = slide.placeholders[1]
        tf = body.text_frame
        
        for point in content['supporting_points']:
            p = tf.add_paragraph()
            p.text = point
            p.font.size = self.bullet_font_size
            p.level = 0
        
        # データの表示
        if 'data' in content:
            data_text = "\n".join([f"{k}: {v}" for k, v in content['data'].items()])
            p = tf.add_paragraph()
            p.text = data_text
            p.font.size = self.bullet_font_size
            p.level = 1
    
    def _create_solution_slide(self, prs, content, design_pattern):
        """解決策スライドの生成"""
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        
        # タイトル
        title = slide.shapes.title
        title.text = content['main_message']
        title.text_frame.paragraphs[0].font.size = self.body_font_size
        title.text_frame.paragraphs[0].font.bold = True
        
        # 本文
        body = slide.placeholders[1]
        tf = body.text_frame
        
        # 解決策
        p = tf.add_paragraph()
        p.text = "解決策:"
        p.font.size = self.bullet_font_size
        p.font.bold = True
        
        for solution in content['solution_points']:
            p = tf.add_paragraph()
            p.text = solution
            p.font.size = self.bullet_font_size
            p.level = 0
        
        # 便益
        p = tf.add_paragraph()
        p.text = "\n期待される便益:"
        p.font.size = self.bullet_font_size
        p.font.bold = True
        
        for benefit in content['benefits']:
            p = tf.add_paragraph()
            p.text = benefit
            p.font.size = self.bullet_font_size
            p.level = 0
    
    def _create_chart_slide(self, prs, content, design_pattern):
        """チャートスライドの生成"""
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        
        # タイトル
        title = slide.shapes.title
        title.text = content['main_message']
        title.text_frame.paragraphs[0].font.size = self.body_font_size
        title.text_frame.paragraphs[0].font.bold = True
        
        # 本文
        body = slide.placeholders[1]
        tf = body.text_frame
        
        # タイムライン
        for phase in content['timeline']:
            p = tf.add_paragraph()
            p.text = f"{phase['phase']} ({phase['duration']})"
            p.font.size = self.bullet_font_size
            p.font.bold = True
            
            for activity in phase['activities']:
                p = tf.add_paragraph()
                p.text = activity
                p.font.size = self.bullet_font_size
                p.level = 0
    
    def _create_financial_slide(self, prs, content, design_pattern):
        """財務スライドの生成"""
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        
        # タイトル
        title = slide.shapes.title
        title.text = content['main_message']
        title.text_frame.paragraphs[0].font.size = self.body_font_size
        title.text_frame.paragraphs[0].font.bold = True
        
        # 本文
        body = slide.placeholders[1]
        tf = body.text_frame
        
        # 投資額
        p = tf.add_paragraph()
        p.text = f"初期投資額: {content['financial_data']['investment']}"
        p.font.size = self.bullet_font_size
        
        # 収益予測
        p = tf.add_paragraph()
        p.text = "\n収益予測:"
        p.font.size = self.bullet_font_size
        p.font.bold = True
        
        for projection in content['financial_data']['revenue_projection']:
            p = tf.add_paragraph()
            p.text = f"Year {projection['year']}: {projection['amount']}"
            p.font.size = self.bullet_font_size
            p.level = 0
        
        # ROI
        p = tf.add_paragraph()
        p.text = f"\n投資回収期間: {content['financial_data']['roi']}"
        p.font.size = self.bullet_font_size
    
    def _create_conclusion_slide(self, prs, content, design_pattern):
        """結論スライドの生成"""
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        
        # タイトル
        title = slide.shapes.title
        title.text = content['main_message']
        title.text_frame.paragraphs[0].font.size = self.body_font_size
        title.text_frame.paragraphs[0].font.bold = True
        
        # 本文
        body = slide.placeholders[1]
        tf = body.text_frame
        
        # アクションアイテム
        p = tf.add_paragraph()
        p.text = "次のステップ:"
        p.font.size = self.bullet_font_size
        p.font.bold = True
        
        for item in content['action_items']:
            p = tf.add_paragraph()
            p.text = item
            p.font.size = self.bullet_font_size
            p.level = 0
        
        # タイムライン
        p = tf.add_paragraph()
        p.text = f"\n実行タイムライン: {content['timeline']}"
        p.font.size = self.bullet_font_size
        
        # 連絡先
        p = tf.add_paragraph()
        p.text = f"\n連絡先: {content['contact_info']}"
        p.font.size = self.bullet_font_size
    
    def save_presentation(self, prs, filename):
        """プレゼンテーションを保存"""
        prs.save(filename) 